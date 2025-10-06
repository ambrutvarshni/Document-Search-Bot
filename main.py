import os
import pickle
from typing import List, Tuple

import docx
import faiss
import gradio as gr
import numpy as np
import pandas as pd
import PyPDF2
from pptx import Presentation
from sentence_transformers import SentenceTransformer

# ---------------- Configuration ----------------
UPLOAD_FOLDER = "uploaded_docs"
EMBEDDINGS_FILE = "embeddings.pkl"
MAX_FILE_SIZE = 2 * 1024 * 1024  # 2MB
CHUNK_SIZE = 200  # smaller chunks for sharper answers
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# ---------------- Global Storage ----------------
document_store = {}  # file_name -> list of embeddings
document_texts = {}  # file_name -> list of text chunks


# ---------------- Persistence ----------------
def save_embeddings():
    with open(EMBEDDINGS_FILE, "wb") as f:
        pickle.dump({"store": document_store, "texts": document_texts}, f)


def load_embeddings():
    global document_store, document_texts
    if os.path.exists(EMBEDDINGS_FILE):
        with open(EMBEDDINGS_FILE, "rb") as f:
            data = pickle.load(f)
            document_store = data.get("store", {})
            document_texts = data.get("texts", {})


# ---------------- Helpers ----------------
def ensure_upload_folder():
    if not os.path.exists(UPLOAD_FOLDER):
        os.makedirs(UPLOAD_FOLDER)


def split_text_into_chunks(text: str, chunk_size: int = CHUNK_SIZE) -> List[str]:
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunks.append(" ".join(words[i : i + chunk_size]))
    return chunks


def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text


def extract_text_from_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_text_from_excel(file_path: str) -> str:
    df = pd.read_excel(file_path)
    return df.to_string()


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def get_embedding(text: str) -> List[float]:
    return embedding_model.encode(text).tolist()


def process_document(file_path: str, file_name: str):
    ext = os.path.splitext(file_name.lower())[1]
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext == ".docx":
        text = extract_text_from_docx(file_path)
    elif ext in [".xlsx", ".xls"]:
        text = extract_text_from_excel(file_path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    chunks = split_text_into_chunks(text)
    embeddings = [get_embedding(chunk) for chunk in chunks]

    document_store[file_name] = embeddings
    document_texts[file_name] = chunks


# ---------------- Document Operations ----------------
def handle_file_upload(files):
    if not files:
        return "No files uploaded."
    uploaded_files = []
    for file in files:
        if os.path.getsize(file.name) > MAX_FILE_SIZE:
            continue
        filename = os.path.basename(file.name)
        save_path = os.path.join(UPLOAD_FOLDER, filename)
        with open(file.name, "rb") as src, open(save_path, "wb") as dst:
            dst.write(src.read())
        try:
            process_document(save_path, filename)
            uploaded_files.append(filename)
        except Exception as e:
            print(f"Error processing {filename}: {e}")
            continue
    save_embeddings()
    return f"Uploaded and processed: {', '.join(uploaded_files)}"


def delete_document(filename: str):
    if isinstance(filename, list):  # gradio sometimes sends lists
        filename = filename[0]
    if filename in document_store:
        del document_store[filename]
        del document_texts[filename]
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(file_path):
            os.remove(file_path)
        save_embeddings()
        return f"Deleted {filename}"
    return f"File {filename} not found"


def search_documents(query: str, top_k: int = 1) -> Tuple[str, str]:
    if not document_store:
        return "No documents uploaded yet.", ""

    query_embedding = get_embedding(query)
    all_embeddings = []
    mapping = []  # (file_name, chunk_index)

    for fname, chunks in document_store.items():
        for idx, emb in enumerate(chunks):
            all_embeddings.append(emb)
            mapping.append((fname, idx))

    if not all_embeddings:
        return "No relevant documents found.", ""

    embeddings_array = np.array(all_embeddings).astype("float32")
    index = faiss.IndexFlatL2(len(query_embedding))
    index.add(embeddings_array)

    D, I = index.search(np.array([query_embedding], dtype="float32"), top_k)
    best_idx = I[0][0]
    fname, chunk_idx = mapping[best_idx]
    return fname, document_texts[fname][chunk_idx]


def query_documents(query: str) -> str:
    fname, chunk = search_documents(query)
    if not chunk:
        return fname
    return f"Closest match from {fname}:\n\n{chunk}"


# ---------------- Gradio UI ----------------
# ---------------- Gradio UI ----------------
def create_ui():
    with gr.Blocks() as app:
        gr.Markdown(
            "<h1 style='text-align:center; color:#4A90E2;'>📑 Document Q&A System</h1>"
        )

        # --- Login Page ---
        with gr.Group() as login_col:
            gr.Markdown("### 🔐 Login")
            role_dropdown = gr.Dropdown(choices=["Admin", "User"], label="Select Role")
            login_btn = gr.Button("Login")
            login_msg = gr.Textbox(label="Login Status", interactive=False)

        # --- Admin Page ---
        with gr.Group(visible=False) as admin_col:
            gr.Markdown("### 🛠️ Admin Dashboard")
            back_admin_btn = gr.Button("🔙 Back to Login")
            file_upload = gr.File(
                file_count="multiple", file_types=[".pdf", ".docx", ".pptx", ".xlsx"]
            )
            upload_btn = gr.Button("Upload Documents")
            upload_output = gr.Textbox(label="Upload Status")
            doc_list = gr.Dropdown(
                choices=list(document_store.keys()), label="Select Document to Delete"
            )
            delete_btn = gr.Button("Delete Document")
            delete_output = gr.Textbox(label="Delete Status")

        # --- User Page ---
        with gr.Group(visible=False) as user_col:
            gr.Markdown("### 🙋 User Dashboard")
            back_user_btn = gr.Button("🔙 Back to Login")
            query_input = gr.Textbox(
                label="Ask a Question", placeholder="Type your query here..."
            )
            query_btn = gr.Button("Search")
            response_output = gr.Textbox(label="Answer", lines=8)

        # --- Login Logic ---
        def handle_login(role):
            if role == "Admin":
                return (
                    "Logged in as Admin",
                    gr.update(visible=False),
                    gr.update(visible=True),
                    gr.update(visible=False),
                )
            elif role == "User":
                return (
                    "Logged in as User",
                    gr.update(visible=False),
                    gr.update(visible=False),
                    gr.update(visible=True),
                )
            else:
                return (
                    "Select a role",
                    gr.update(visible=True),
                    gr.update(visible=False),
                    gr.update(visible=False),
                )

        login_btn.click(
            handle_login,
            inputs=[role_dropdown],
            outputs=[login_msg, login_col, admin_col, user_col],
        )

        # --- Back Buttons ---
        back_admin_btn.click(
            lambda: (
                gr.update(visible=True),
                gr.update(visible=False),
                gr.update(visible=False),
            ),
            outputs=[login_col, admin_col, user_col],
        )
        back_user_btn.click(
            lambda: (
                gr.update(visible=True),
                gr.update(visible=False),
                gr.update(visible=False),
            ),
            outputs=[login_col, admin_col, user_col],
        )

        # --- Admin Actions ---
        def upload_and_refresh(files):
            msg = handle_file_upload(files)
            return msg, list(document_store.keys())

        upload_btn.click(
            upload_and_refresh, inputs=[file_upload], outputs=[upload_output, doc_list]
        )

        def delete_and_refresh(filename):
            msg = delete_document(filename)
            return msg, list(document_store.keys())

        delete_btn.click(
            delete_and_refresh, inputs=[doc_list], outputs=[delete_output, doc_list]
        )

        # --- User Actions ---
        def precise_query(query: str):
            fname, chunk = search_documents(query, top_k=3)
            if not chunk:
                return fname
            # Just return first 2-3 sentences for precision
            sentences = chunk.split(". ")
            snippet = ". ".join(sentences[:3])
            return f"Closest match from {fname}:\n\n{snippet}..."

        query_btn.click(precise_query, inputs=[query_input], outputs=[response_output])

    return app


# ---------------- Main ----------------
if __name__ == "__main__":
    ensure_upload_folder()
    load_embeddings()
    app = create_ui()
    app.launch()
