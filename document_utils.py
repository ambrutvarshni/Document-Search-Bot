import os
import pickle
from typing import List, Tuple

import docx
import faiss
import numpy as np
import pandas as pd
import PyPDF2
from pptx import Presentation
from sentence_transformers import SentenceTransformer

# ---------------- Constants ----------------
UPLOAD_FOLDER = "uploaded_docs"
EMBEDDINGS_FILE = "embeddings.pkl"
MAX_FILE_SIZE = 2 * 1024 * 1024  # 2MB
CHUNK_SIZE = 200  # words per chunk

# ---------------- Globals ----------------
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
document_store = {}  # file_name -> list of embeddings
document_texts = {}  # file_name -> list of chunks


# ---------------- Persistence ----------------
def save_embeddings():
    with open(EMBEDDINGS_FILE, "wb") as f:
        pickle.dump({"store": document_store, "texts": document_texts}, f)


def load_embeddings():
    global document_store, document_texts
    if os.path.exists(EMBEDDINGS_FILE):
        with open(EMBEDDINGS_FILE, "rb") as f:
            data = pickle.load(f)
            document_store = data.get("store", {})
            document_texts = data.get("texts", {})


# ---------------- Utilities ----------------
def ensure_upload_folder():
    if not os.path.exists(UPLOAD_FOLDER):
        os.makedirs(UPLOAD_FOLDER)


def split_text_into_chunks(text: str, chunk_size: int = CHUNK_SIZE) -> List[str]:
    words = text.split()
    return [
        " ".join(words[i : i + chunk_size]) for i in range(0, len(words), chunk_size)
    ]


def get_embedding(text: str) -> List[float]:
    return embedding_model.encode(text).tolist()


# ---------------- Extract Text ----------------
def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text


def extract_text_from_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_text_from_excel(file_path: str) -> str:
    df = pd.read_excel(file_path)
    return df.to_string()


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# ---------------- Document Processing ----------------
def process_document(file_path: str, file_name: str):
    ext = os.path.splitext(file_name.lower())[1]
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext == ".docx":
        text = extract_text_from_docx(file_path)
    elif ext in [".xlsx", ".xls"]:
        text = extract_text_from_excel(file_path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    chunks = split_text_into_chunks(text)
    embeddings = [get_embedding(chunk) for chunk in chunks]

    document_store[file_name] = embeddings
    document_texts[file_name] = chunks


# ---------------- Upload/Delete ----------------
def handle_file_upload(files):
    if not files:
        return "No files uploaded.", list(document_store.keys())
    uploaded_files = []
    for file in files:
        if os.path.getsize(file.name) > MAX_FILE_SIZE:
            continue
        filename = os.path.basename(file.name)
        save_path = os.path.join(UPLOAD_FOLDER, filename)
        with open(file.name, "rb") as src, open(save_path, "wb") as dst:
            dst.write(src.read())
        try:
            process_document(save_path, filename)
            uploaded_files.append(filename)
        except Exception as e:
            print(f"Error processing {filename}: {e}")
            continue
    save_embeddings()
    return f"Uploaded: {', '.join(uploaded_files)}", list(document_store.keys())


def delete_document(filename):
    if isinstance(filename, list):
        filename = filename[0] if filename else ""
    if not filename:
        return "No file selected", list(document_store.keys())
    if filename in document_store:
        del document_store[filename]
        del document_texts[filename]
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(file_path):
            os.remove(file_path)
        save_embeddings()
        return f"Deleted {filename}", list(document_store.keys())
    return f"File {filename} not found", list(document_store.keys())


# ---------------- Search ----------------
def search_documents_topk(query: str, top_k: int = 3) -> Tuple[str, List[str]]:
    """
    Returns the filename and a list of top-k text chunks most similar to the query.
    """
    if not document_store:
        return "No documents uploaded yet.", []

    query_emb = get_embedding(query)
    all_embeddings = []
    mapping = []

    for fname, chunks in document_store.items():
        for idx, emb in enumerate(chunks):
            all_embeddings.append(emb)
            mapping.append((fname, idx))

    embeddings_array = np.array(all_embeddings).astype("float32")
    index = faiss.IndexFlatL2(len(query_emb))
    index.add(embeddings_array)
    D, I = index.search(np.array([query_emb], dtype="float32"), top_k)

    top_chunks = []
    for idx in I[0]:
        fname, chunk_idx = mapping[idx]
        top_chunks.append(document_texts[fname][chunk_idx])

    if not top_chunks:
        return "No relevant content found.", []

    return fname, top_chunks
